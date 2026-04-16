import os
import uuid
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
# Point pytesseract to the installed Tesseract binary (adjust path if different)
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
from io import BytesIO

from embedder import model
from chunking import hybrid_chunk_document, get_avg_tokens_per_chunk, chunk_text_token_aware
from vector_store import add_document_chunks as add_vector_chunks, search_context
from db import save_document_metadata, add_document_chunks as add_db_chunks
from keyword_search import keyword_search  # <- required for hybrid retrieval
from multimodal_processor import multimodal_processor

from transformers import BlipProcessor, BlipForConditionalGeneration
from bs4 import BeautifulSoup
import docx
from pptx import Presentation
from ebooklib import epub

# Load BLIP captioning model globally
try:
    blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base", use_fast=True)
    blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
except Exception as e:
    print(f"Warning: BLIP model could not be loaded: {e}")
    blip_processor = None
    blip_model = None


# --- Synchronous Processing for Windows Stability ---

def process_image_locally(page_index, img_index, image_bytes):
    """Processes a single image within a PDF page."""
    try:
        image = Image.open(BytesIO(image_bytes)).convert("RGB")
        ocr_text = pytesseract.image_to_string(image).strip()

        caption = ""
        if blip_model and blip_processor:
            inputs = blip_processor(image, return_tensors="pt")
            caption_ids = blip_model.generate(**inputs, max_new_tokens=30)
            caption = blip_processor.decode(caption_ids[0], skip_special_tokens=True)

        result = f"\n[Image Page {page_index + 1}, Image {img_index + 1}]\n"
        if caption:
            result += f"Caption: {caption}\n"
        if ocr_text:
            result += f"OCR: {ocr_text}"
        return result.strip()
    except Exception as e:
        return f" OCR/Captioning failed on page {page_index + 1}, image {img_index + 1}: {e}"


def extract_image_text(pdf_path):
    """Extracts text from images embedded in a PDF."""
    doc = fitz.open(pdf_path)
    combined_results = []

    for page_index in range(len(doc)):
        page = doc.load_page(page_index)
        for img_index, img in enumerate(page.get_images(full=True)):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                res = process_image_locally(page_index, img_index, image_bytes)
                combined_results.append(res)
            except Exception as e:
                print(f" Image extract failed on page {page_index + 1}, image {img_index + 1}: {e}")

    doc.close()
    return "\n\n".join(combined_results)


# --- Format-specific text extractors ---
def extract_text_from_docx(path):
    return "\n".join(p.text for p in docx.Document(path).paragraphs if p.text.strip())


def extract_text_from_pptx(path):
    prs = Presentation(path)
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))


def extract_text_from_html(path):
    with open(path, encoding="utf-8") as f:
        return BeautifulSoup(f, "html.parser").get_text(separator="\n")


def extract_text_from_epub(path):
    book = epub.read_epub(path)
    texts = []
    for item in book.get_items():
        if item.get_type() == epub.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            texts.append(soup.get_text(separator="\n"))
    return "\n".join(texts)


# --- Document Ingestion Pipeline ---

def ingest_multimodal_document(file_path):
    """Ingest an image / video / audio file via the multimodal processor."""
    print(f"\n🎬 [MULTIMODAL] Processing: {file_path}")

    try:
        media_result = multimodal_processor.process_media_file(file_path)
        if "error" in media_result:
            print(f"❌ Media processing failed: {media_result['error']}")
            return None, 0, None

        searchable_content = multimodal_processor.extract_media_content_for_chunking(file_path)
        doc_id = str(uuid.uuid4())[:8]
        chunks = chunk_text_token_aware(searchable_content, max_tokens=300, overlap=32)
        
        # Embed synchronously
        from embedder import model
        embeddings = [model.encode([c])[0].tolist() for c in chunks]
        add_vector_chunks(doc_id, chunks, embeddings)
        add_db_chunks(doc_id, chunks, embeddings)
        avg_tokens = get_avg_tokens_per_chunk(chunks)
        save_document_metadata(doc_id, os.path.basename(file_path), avg_tokens)

        print(f"✅ [MULTIMODAL] Indexed {len(chunks)} chunks for {os.path.basename(file_path)}")
        return len(chunks), avg_tokens, doc_id

    except Exception as e:
        print(f"❌ [MULTIMODAL] Ingestion failed for {file_path}: {e}")
        return None, 0, None


def ingest_document(file_path):
    """Main entry point for document ingestion. Handles all formats synchronously."""
    ext = os.path.splitext(file_path)[1].lower()
    combined_text = ""
    doc_id = str(uuid.uuid4())[:8]
    
    try:
        # Detect file type and route accordingly
        if ext in multimodal_processor.supported_image_formats + \
           multimodal_processor.supported_video_formats + \
           multimodal_processor.supported_audio_formats:
            return ingest_multimodal_document(file_path)

        # Process based on file extension
        print(f"🔍 [INGEST] Parsing format: {ext}")
        
        if ext == '.pdf':
            from embedder import extract_text as extract_text_pdf
            combined_text = extract_text_pdf(file_path)
            
            # --- Robust OCR Fallback for Scanned PDFs (Self-Contained Rendering) ---
            if len(combined_text.strip()) < 100:
                print(" ⚠️ Scanned PDF detected (low content). Running Full-Page OCR...")
                doc = fitz.open(file_path)
                for i, page in enumerate(doc):
                    print(f"   OCR Processing Page {i+1}...")
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2)) # 144 DPI sufficient
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    page_text = pytesseract.image_to_string(img).strip()
                    combined_text += f"\n\n[OCR Page {i+1}]\n" + page_text
                doc.close()
                print(f" ✅ OCR recovered {len(combined_text)} characters.")
            
            # Add text from internal images (via BLIP captioning)
            combined_text += "\n\n" + extract_image_text(file_path)

            # --- Local Visual Analysis: Tables, Charts, Graphs ---
            # This appends table data and image/chart descriptions to the text.
            # Uses pdfplumber (tables) + Tesseract (chart labels) + BLIP (captions).
            # No API required — runs fully locally.
            try:
                from local_visual_analyzer import analyze_document_visuals
                visual_content = analyze_document_visuals(file_path)
                if visual_content:
                    combined_text += "\n\n" + visual_content
                    print(f" ✅ Visual content (tables/charts) extracted and appended.")
            except Exception as ve:
                print(f" ⚠️ Visual analysis skipped: {ve}")

            
        elif ext == '.docx':
            combined_text = extract_text_from_docx(file_path)
        elif ext == '.pptx':
            combined_text = extract_text_from_pptx(file_path)
        elif ext == '.html':
            combined_text = extract_text_from_html(file_path)
        elif ext == '.epub':
            combined_text = extract_text_from_epub(file_path)
        elif ext == '.txt':
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                combined_text = f.read()
        else:
            raise ValueError(f"Unsupported file type: {ext}")

        if not combined_text.strip():
            print(f"⚠️ [INGEST] No text content found for {file_path}")
            return 0, 0, doc_id

        # Chunk and embed (Synchronously for stability)
        chunks = chunk_text_token_aware(combined_text, max_tokens=300, overlap=32)
        
        from embedder import model
        if model is not None:
             embeddings = [model.encode([c])[0].tolist() for c in chunks]
        else:
             embeddings = [[0.0]*384 for _ in chunks]

        add_vector_chunks(doc_id, chunks, embeddings)
        add_db_chunks(doc_id, chunks, embeddings)
        avg_tokens = get_avg_tokens_per_chunk(chunks)
        save_document_metadata(doc_id, os.path.basename(file_path), avg_tokens)

        print(f"✅ [INGEST] Perfectly indexed {len(chunks)} chunks for {os.path.basename(file_path)}")
        return len(chunks), avg_tokens, doc_id
        
    except Exception as e:
        print(f"[❌] Ingestion failed for {file_path}: {e}")
        raise RuntimeError(f"Document ingestion failed: {e}")


# --- Hybrid Retrieval for RAG ---

def retrieve_context(query, top_k=5):
    """Retrieves context using Semantic and Keyword search."""
    print(f"\n[RAG] Retrieving hybrid context for query: {query}")
    
    dense_results = []
    keyword_results = []

    try:
        dense_results = search_context(query, top_k=top_k)
    except Exception as e:
        print(f"Dense search failed: {e}")

    try:
        keyword_results = keyword_search(query, top_k=top_k)
    except Exception as e:
        print(f"Keyword search failed: {e}")

    # Combine and deduplicate
    seen = set()
    combined = []
    for r in dense_results + keyword_results:
        text = r[0] if isinstance(r, tuple) else r
        if text not in seen:
            combined.append(text)
            seen.add(text)

    return combined[:top_k]


def get_layered_context(query, top_k=5):
    """Implements a Layered Retrieval Strategy with fallback to basic policy."""
    context = retrieve_context(query, top_k=top_k)
    
    if not context:
        print(f"⚠️ [LayeredRAG] Falling back to basic_hr_policy.md")
        try:
            with open("basic_hr_policy.md", "r", encoding="utf-8") as f:
                return [f.read()]
        except FileNotFoundError:
            return []
            
    return context


def process_rag_query(query):
    """End-to-end RAG query processing."""
    from gemini_llm import query_gemini
    context = get_layered_context(query, top_k=5)
    
    if not context:
        return "I don't have specific information about that query in the company policy documents."
    
    context_text = "\n\n".join(context)
    try:
        response = query_gemini([context_text], query)
        return response
    except Exception as e:
        return f"I encountered an error while processing your query. ({str(e)})"