import os
import uuid
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
# Point pytesseract to the installed Tesseract binary (adjust path if different)
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
from io import BytesIO
import multiprocessing

from embedder import model
from chunking import hybrid_chunk_document, get_avg_tokens_per_chunk, chunk_text_token_aware
from vector_store import add_document_chunks, search_context
from db import save_document_metadata
from keyword_search import keyword_search  # <- required for hybrid retrieval
from multimodal_processor import multimodal_processor

from transformers import BlipProcessor, BlipForConditionalGeneration
from bs4 import BeautifulSoup
import docx
from pptx import Presentation
from ebooklib import epub

# Windows multiprocessing safety
multiprocessing.freeze_support()

# Load BLIP captioning model
blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base", use_fast=True)
blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")


# --- Multiprocessing-safe embedding worker ---
def embedding_worker(chunk):
    from embedder import model  # re-import inside subprocess
    return model.encode([chunk])[0]


def embed_chunks_parallel(chunks):
    with multiprocessing.Pool(processes=min(4, multiprocessing.cpu_count())) as pool:
        embeddings = pool.map(embedding_worker, chunks)
    return embeddings


# --- Image OCR & Captioning task ---
def process_image_task(task):
    page_index, img_index, image_bytes = task
    try:
        image = Image.open(BytesIO(image_bytes)).convert("RGB")
        ocr_text = pytesseract.image_to_string(image).strip()

        inputs = blip_processor(image, return_tensors="pt")
        caption_ids = blip_model.generate(**inputs, max_new_tokens=30)
        caption = blip_processor.decode(caption_ids[0], skip_special_tokens=True)

        result = f"[Image Page {page_index + 1}, Image {img_index + 1}]\n"
        if caption:
            result += f"Caption: {caption}\n"
        if ocr_text:
            result += f"OCR: {ocr_text}"
        return result.strip()
    except Exception as e:
        return f" OCR/Captioning failed on page {page_index + 1}, image {img_index + 1}: {e}"


def extract_image_text(pdf_path):
    doc = fitz.open(pdf_path)
    image_tasks = []

    for page_index in range(len(doc)):
        page = doc.load_page(page_index)
        for img_index, img in enumerate(page.get_images(full=True)):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_tasks.append((page_index, img_index, image_bytes))
            except Exception as e:
                print(f" Image extract failed on page {page_index + 1}, image {img_index + 1}: {e}")

    with multiprocessing.Pool(processes=min(4, multiprocessing.cpu_count())) as pool:
        results = pool.map(process_image_task, image_tasks)

    return "\n\n".join(results)


# --- Format-specific text extractors ---
def extract_text_from_docx(path):
    return "\n".join(p.text for p in docx.Document(path).paragraphs if p.text.strip())


def extract_text_from_pptx(path):
    prs = Presentation(path)
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))


def extract_text_from_html(path):
    with open(path, encoding="utf-8") as f:
        return BeautifulSoup(f, "html.parser").get_text(separator="\n")


def extract_text_from_epub(path):
    book = epub.read_epub(path)
    texts = []
    for item in book.get_items():
        if item.get_type() == epub.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            texts.append(soup.get_text(separator="\n"))
    return "\n".join(texts)


def ingest_multimodal_document(file_path):
    """Ingest an image / video / audio file via the multimodal processor."""
    print(f"\n🎬 [MULTIMODAL] Processing: {file_path}")

    try:
        # Let the multimodal processor analyse the media file
        media_result = multimodal_processor.process_media_file(file_path)

        if "error" in media_result:
            print(f"❌ Media processing failed: {media_result['error']}")
            return None, 0, None

        # Turn the structured result into a single text blob we can chunk
        searchable_content = multimodal_processor.extract_media_content_for_chunking(file_path)

        # For multimodal we already have text, so chunk it directly by tokens
        doc_id = str(uuid.uuid4())[:8]
        chunks = chunk_text_token_aware(searchable_content, max_tokens=300, overlap=32)
        embeddings = embed_chunks_parallel(chunks)
        add_document_chunks(doc_id, chunks, embeddings)
        avg_tokens = get_avg_tokens_per_chunk(chunks)
        save_document_metadata(doc_id, os.path.basename(file_path), avg_tokens)

        print(f"✅ [MULTIMODAL] Indexed {len(chunks)} chunks for {os.path.basename(file_path)}")
        return len(chunks), avg_tokens, doc_id

    except Exception as e:
        print(f"❌ [MULTIMODAL] Ingestion failed for {file_path}: {e}")
        return None, 0, None


# --- Document Ingestion Pipeline ---
def ingest_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    combined_text = ""
    try:
        # Detect file type and route accordingly
        file_ext = os.path.splitext(file_path)[1].lower()

        # Check if it's a media file first
        if file_ext in multimodal_processor.supported_image_formats + \
           multimodal_processor.supported_video_formats + \
           multimodal_processor.supported_audio_formats:
            return ingest_multimodal_document(file_path)

        # Otherwise, treat as a regular text document
        combined_text = ""
        doc_id = str(uuid.uuid4())[:8]

        # Process based on file extension
        if file_ext == '.pdf':
            print(" Parsing PDF...")
            from embedder import extract_text as extract_text_pdf
            combined_text += extract_text_pdf(file_path)
            combined_text += "\n\n" + extract_image_text(file_path)
        elif file_ext == '.docx':
            print(" Parsing DOCX...")
            combined_text += extract_text_from_docx(file_path)
        elif file_ext == '.pptx':
            print(" Parsing PPTX...")
            combined_text += extract_text_from_pptx(file_path)
        elif file_ext == '.html':
            print(" Parsing HTML...")
            combined_text += extract_text_from_html(file_path)
        elif file_ext == '.epub':
            print(" Parsing EPUB...")
            combined_text += extract_text_from_epub(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")

        # Chunk and embed (we already have full text, so use token-aware text chunker)
        chunks = chunk_text_token_aware(combined_text, max_tokens=300, overlap=32)
        embeddings = embed_chunks_parallel(chunks)
        add_document_chunks(doc_id, chunks, embeddings)
        avg_tokens = get_avg_tokens_per_chunk(chunks)
        save_document_metadata(doc_id, os.path.basename(file_path), avg_tokens)

        return len(chunks), avg_tokens, doc_id
    except Exception as e:
        print(f"[❌] Ingestion failed for {file_path}: {e}")
        raise RuntimeError(f"Document ingestion failed: {e}")


# --- Hybrid Retrieval for RAG ---
def retrieve_context(query, top_k=5):
    print(f"\n[RAG] Retrieving hybrid context for query: {query}")
    
    dense_results = []
    keyword_results = []

    try:
        dense_results = search_context(query, top_k=top_k)
        print(f"Dense results: {len(dense_results)}")
    except Exception as e:
        print(f"Dense search failed for query '{query}': {e}")

    try:
        keyword_results = keyword_search(query, top_k=top_k)
        print(f"Keyword results: {len(keyword_results)}")
    except Exception as e:
        print(f"Keyword search failed for query '{query}': {e}")

    # Combine and deduplicate
    seen = set()
    combined = []
    for r in dense_results + keyword_results:
        text = r[0] if isinstance(r, tuple) else r
        if text not in seen:
            combined.append(text)
            seen.add(text)

    if not combined:
        print(f"⚠️ [RAG] No context found for query: {query}")
    else:
        print(f"[RAG] Returning {len(combined[:top_k])} hybrid context chunks.")

    return combined[:top_k]


def get_layered_context(query, top_k=5):
    """
    Implements a Layered Retrieval Strategy:
    1. Vector Database (Chroma) + Keyword Search
    2. Fallback to basic_hr_policy.md if nothing is found
    """
    # 1. Standard Hybrid Retrieval
    context = retrieve_context(query, top_k=top_k)
    
    # 2. Fallback to basic document
    if not context:
        print(f"⚠️ [LayeredRAG] Vector DB returned empty. Falling back to basic_hr_policy.md")
        try:
            with open("basic_hr_policy.md", "r", encoding="utf-8") as f:
                fallback_text = f.read()
                return [fallback_text]
        except FileNotFoundError:
            print("⚠️ [LayeredRAG] Fallback document not found.")
            return []
            
    return context


def process_rag_query(query):
    """
    Process a RAG query by retrieving context and generating a response.
    """
    import gemini_llm
    
    # Use Layered Retrieval
    context = get_layered_context(query, top_k=5)
    
    # If no context found, return a fallback response
    if not context:
        return "I don't have specific information about that query in the company policy documents."
    
    # Combine context into a single string
    context_text = "\n\n".join(context)
    
    # Generate response using LLM
    try:
        response = gemini_llm.query_gemini([context_text], query)
        return response
    except Exception as e:
        print(f"Error processing RAG query: {e}")
        return "I encountered an error while processing your query. Please try again or contact HR."